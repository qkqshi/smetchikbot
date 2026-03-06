import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def generate_kp_pptx(project_name, items_data, furniture_data=None, photo_path=None, phone_number=None):
    """
    Генерирует PPTX презентацию КП
    
    Args:
        project_name: название проекта
        items_data: список словарей с данными об изделиях (материалы)
        furniture_data: список словарей с данными о мебели
        photo_path: путь к фото проекта
        phone_number: номер телефона для обложки (опционально)
    
    Returns:
        путь к созданному PPTX файлу
    """
    print(f"[PPTX] Начало генерации для проекта: {project_name}")
    print(f"[PPTX] Фото проекта: {photo_path}, существует: {os.path.exists(photo_path) if photo_path else False}")
    if phone_number:
        print(f"[PPTX] Номер телефона для обложки: {phone_number}")
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # СЛАЙД 1: Обложка
    print("[PPTX] Создание слайда обложки...")
    cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Ищем обложку в разных форматах и местах
    cover_names = ['cover.jpg', 'cover.png', 'cover.jpeg']
    base_paths = [
        'templates',
        '/opt/smetchikbot/templates',
        os.path.join(os.path.dirname(__file__), '..', 'templates')
    ]
    
    possible_cover_paths = []
    for base in base_paths:
        for name in cover_names:
            possible_cover_paths.append(os.path.join(base, name))
    
    cover_added = False
    for cover_path in possible_cover_paths:
        if os.path.exists(cover_path):
            try:
                print(f"[PPTX] Найдена обложка: {cover_path}")
                # Проверяем валидность обложки
                from PIL import Image as PILImage
                img = PILImage.open(cover_path)
                img_width, img_height = img.size
                print(f"[PPTX] Размер обложки: {img_width}x{img_height}")
                
                # Растягиваем обложку на весь слайд (без сохранения пропорций для обложки)
                cover_slide.shapes.add_picture(cover_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                print("[PPTX] Обложка добавлена успешно")
                
                # Если указан номер телефона, добавляем его в пустой прямоугольник
                if phone_number:
                    print(f"[PPTX] Добавление номера телефона на обложку: {phone_number}")
                    
                    # Текстовый блок с номером (в позиции пустого прямоугольника внизу слева)
                    left = Inches(0.2)
                    top = prs.slide_height - Inches(0.8)
                    width = Inches(3)
                    height = Inches(0.5)
                    
                    phone_box = cover_slide.shapes.add_textbox(left, top, width, height)
                    phone_frame = phone_box.text_frame
                    phone_frame.text = f"тел. {phone_number}"
                    phone_frame.vertical_anchor = 1  # Вертикальное выравнивание по центру
                    phone_para = phone_frame.paragraphs[0]
                    phone_para.font.size = Pt(14)
                    phone_para.font.bold = True
                    phone_para.font.color.rgb = RGBColor(255, 255, 255)  # Белый цвет
                    phone_para.alignment = PP_ALIGN.LEFT
                    
                    # Добавляем отступ слева
                    phone_frame.margin_left = Inches(0.2)
                
                cover_added = True
                break
            except Exception as e:
                print(f"Ошибка добавления обложки {cover_path}: {e}")
                import traceback
                traceback.print_exc()
    
    if not cover_added:
        print(f"[PPTX] Обложка не найдена. Проверенные пути:")
        for path in possible_cover_paths:
            print(f"  - {path} (существует: {os.path.exists(path)})")
    
    # СЛАЙДЫ С МАТЕРИАЛАМИ (по одному слайду на каждую позицию)
    print(f"[PPTX] Создание слайдов для материалов, позиций: {len(items_data)}")
    for idx, item_data in enumerate(items_data):
        print(f"[PPTX] Обработка материала {idx+1}/{len(items_data)}: {item_data['name']}")
        material_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Название материала
        material_title_box = material_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        material_title_frame = material_title_box.text_frame
        material_title_frame.text = item_data['name']
        material_title_para = material_title_frame.paragraphs[0]
        material_title_para.font.size = Pt(20)
        material_title_para.font.bold = True
        material_title_para.font.color.rgb = RGBColor(0, 0, 0)
        material_title_para.alignment = PP_ALIGN.CENTER
        
        y_pos = 1.0
        
        # Добавляем фото материала (если есть) с сохранением пропорций
        item_photo = item_data.get('image')  # Индивидуальное фото материала
        if not item_photo or not os.path.exists(item_photo):
            item_photo = photo_path  # Fallback на общее фото проекта
        
        if item_photo and os.path.exists(item_photo):
            try:
                print(f"[PPTX] Добавление фото на слайд материала: {item_photo}")
                from PIL import Image as PILImage
                img = PILImage.open(item_photo)
                img_width, img_height = img.size
                print(f"[PPTX] Фото валидно, размер: {img_width}x{img_height}")
                
                # Максимальные размеры
                max_width = Inches(5)
                max_height = Inches(3)
                
                # Вычисляем масштаб с сохранением пропорций
                width_ratio = max_width / img_width
                height_ratio = max_height / img_height
                scale = min(width_ratio, height_ratio)
                
                # Новые размеры
                new_width = img_width * scale
                new_height = img_height * scale
                
                # Центрируем по горизонтали
                left = (prs.slide_width - new_width) / 2
                top = Inches(y_pos)
                
                material_slide.shapes.add_picture(item_photo, left, top, width=new_width, height=new_height)
                print(f"[PPTX] Фото добавлено на слайд материала: {new_width/Inches(1):.2f}x{new_height/Inches(1):.2f} дюймов")
                y_pos += (new_height / Inches(1)) + 0.3
            except Exception as e:
                print(f"Ошибка добавления изображения на слайд материала: {e}")
                import traceback
                traceback.print_exc()
        
        # Детали материала
        details = [
            f"Каркас из {item_data.get('body_material', 'ЛДСП')}",
            f"Фасады {item_data.get('facade_description', 'МДФ')}",
        ]
        
        if item_data.get('additional_info'):
            details.extend(item_data['additional_info'])
        
        info_box = material_slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(2))
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        
        info_text = "\n".join(details)
        quantity = item_data.get('quantity', 1)
        info_text += f"\n\nКоличество: {quantity} шт."
        info_text += f"\nИтоговая стоимость: {item_data.get('total_cost', 0):,.1f} руб".replace(',', ' ')
        
        info_frame.text = info_text
        for para in info_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.alignment = PP_ALIGN.CENTER
    
    # СЛАЙДЫ С МЕБЕЛЬЮ (по одному слайду на каждую позицию)
    if furniture_data:
        print(f"[PPTX] Создание слайдов для мебели, позиций: {len(furniture_data)}")
        for idx, furniture_item in enumerate(furniture_data):
            print(f"[PPTX] Обработка мебели {idx+1}/{len(furniture_data)}: {furniture_item['name']}")
            furn_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Название мебели
            furn_title_box = furn_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            furn_title_frame = furn_title_box.text_frame
            furn_title_frame.text = furniture_item['name']
            furn_title_para = furn_title_frame.paragraphs[0]
            furn_title_para.font.size = Pt(20)
            furn_title_para.font.bold = True
            furn_title_para.font.color.rgb = RGBColor(0, 0, 0)
            furn_title_para.alignment = PP_ALIGN.CENTER
            
            # Изображение мебели (если есть) с сохранением пропорций
            y_pos = 1.0
            furn_image = furniture_item.get('image')
            print(f"[PPTX] Изображение мебели: {furn_image}, существует: {os.path.exists(furn_image) if furn_image else False}")
            
            if furn_image and os.path.exists(furn_image):
                try:
                    print(f"[PPTX] Добавление изображения мебели: {furn_image}")
                    # Проверяем валидность изображения
                    from PIL import Image as PILImage
                    img = PILImage.open(furn_image)
                    img_width, img_height = img.size
                    print(f"[PPTX] Изображение мебели валидно, размер: {img_width}x{img_height}")
                    
                    # Максимальные размеры
                    max_width = Inches(6)
                    max_height = Inches(4)
                    
                    # Вычисляем масштаб с сохранением пропорций
                    width_ratio = max_width / img_width
                    height_ratio = max_height / img_height
                    scale = min(width_ratio, height_ratio)
                    
                    # Новые размеры
                    new_width = img_width * scale
                    new_height = img_height * scale
                    
                    # Центрируем по горизонтали
                    left = (prs.slide_width - new_width) / 2
                    top = Inches(y_pos)
                    
                    furn_slide.shapes.add_picture(furn_image, left, top, width=new_width, height=new_height)
                    print(f"[PPTX] Изображение мебели добавлено: {new_width/Inches(1):.2f}x{new_height/Inches(1):.2f} дюймов")
                    y_pos += (new_height / Inches(1)) + 0.5  # Отступ после изображения
                except Exception as e:
                    print(f"Ошибка добавления изображения мебели: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Информация о мебели (под изображением)
            info_box = furn_slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(2))
            info_frame = info_box.text_frame
            info_frame.word_wrap = True
            
            info_text = f"Количество: {furniture_item['quantity']}\n"
            info_text += f"Стоимость за 1 шт: {furniture_item['price_per_unit']} руб\n"
            info_text += f"Итоговая стоимость: {furniture_item['total_price']} руб"
            
            info_frame.text = info_text
            for para in info_frame.paragraphs:
                para.font.size = Pt(18)
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.alignment = PP_ALIGN.CENTER
    
    # СЛАЙД С СОВМЕЩЕННОЙ ТАБЛИЦЕЙ (материалы + мебель)
    print("[PPTX] Создание слайда с совмещенной таблицей...")
    table_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Заголовок
    title_box = table_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = project_name
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Собираем все элементы (материалы + мебель)
    all_items = []
    
    # Добавляем материалы
    for item_data in items_data:
        total = item_data.get('total_cost', 0)
        quantity = item_data.get('quantity', 1)
        all_items.append({
            'name': item_data['name'],
            'quantity': quantity,
            'price_c': round(total, 1),
            'price_d': round(total, 1)
        })
    
    # Добавляем мебель
    if furniture_data:
        for furn in furniture_data:
            quantity = furn['quantity']
            if isinstance(quantity, str):
                quantity = int(quantity)
            
            price_per_unit = furn['price_per_unit']
            if isinstance(price_per_unit, str):
                price_per_unit = float(price_per_unit.replace(' ', '').replace(',', ''))
            
            total_price = furn['total_price']
            if isinstance(total_price, str):
                total_price = float(total_price.replace(' ', '').replace(',', ''))
            
            all_items.append({
                'name': furn['name'],
                'quantity': quantity,
                'price_c': round(price_per_unit, 1),
                'price_d': round(total_price, 1)
            })
    
    # Создаем таблицу (строки: данные + итого, столбцы: 4)
    rows = len(all_items) + 1  # +1 для итого
    cols = 4
    
    # Позиция и размер таблицы
    left = Inches(1.5)
    top = Inches(1.2)
    width = Inches(7)
    height = Inches(0.35) * rows  # Уменьшаем высоту строк для освобождения места
    
    table_shape = table_slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Настройка ширины столбцов
    table.columns[0].width = Inches(3.5)  # Наименование
    table.columns[1].width = Inches(1.0)  # Количество
    table.columns[2].width = Inches(1.25)  # Цена 1
    table.columns[3].width = Inches(1.25)  # Цена 2
    
    # Заполняем данные (без заголовков)
    total_col_c = 0
    total_col_d = 0
    
    for row_idx, item in enumerate(all_items):
        # Наименование
        cell = table.cell(row_idx, 0)
        cell.text = item['name']
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Количество
        cell = table.cell(row_idx, 1)
        cell.text = str(item['quantity'])
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Цена 1
        cell = table.cell(row_idx, 2)
        cell.text = f"{item['price_c']:,.1f}".replace(',', ' ')
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Цена 2
        cell = table.cell(row_idx, 3)
        cell.text = f"{item['price_d']:,.1f}".replace(',', ' ')
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        total_col_c += item['price_c']
        total_col_d += item['price_d']
    
    # Строка "Итого"
    last_row = rows - 1
    cell = table.cell(last_row, 0)
    cell.text = "Итого"
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Пустая ячейка
    cell = table.cell(last_row, 1)
    cell.text = ""
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Итого колонка C
    cell = table.cell(last_row, 2)
    cell.text = f"{total_col_c:,.1f}".replace(',', ' ')
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Итого колонка D (с желтым фоном)
    cell = table.cell(last_row, 3)
    cell.text = f"{total_col_d:,.1f}".replace(',', ' ')
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Желтый
    
    print(f"[PPTX] Совмещенная таблица создана, итого: {total_col_d:,.1f} руб")
    
    # Добавляем примечание под таблицей
    note_top = top + height + Inches(0.2)
    note_box = table_slide.shapes.add_textbox(Inches(1.5), note_top, Inches(7), Inches(1.2))
    note_frame = note_box.text_frame
    note_frame.word_wrap = True
    note_frame.margin_top = Inches(0.05)
    note_frame.margin_bottom = Inches(0.05)
    
    note_text = "Услуги упаковка, доставка, подъем, сборка и монтаж, сбор мусора +10% к стоимости\n"
    note_text += "-3% Скидка за наличные\n"
    note_text += "Срок 45 рабочих дней"
    
    note_frame.text = note_text
    for para in note_frame.paragraphs:
        para.font.size = Pt(10)
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(3)  # Уменьшаем отступ между строками
    
    # ФИНАЛЬНЫЕ СЛАЙДЫ: end1, end2, end3
    print("[PPTX] Добавление финальных слайдов...")
    end_slide_names = ['end1', 'end2', 'end3']
    end_formats = ['.png', '.jpg', '.jpeg']
    
    base_paths = [
        'templates',
        '/opt/smetchikbot/templates',
        os.path.join(os.path.dirname(__file__), '..', 'templates')
    ]
    
    for end_name in end_slide_names:
        end_added = False
        for base in base_paths:
            for fmt in end_formats:
                end_path = os.path.join(base, end_name + fmt)
                if os.path.exists(end_path):
                    try:
                        print(f"[PPTX] Найден финальный слайд: {end_path}")
                        from PIL import Image as PILImage
                        img = PILImage.open(end_path)
                        img_width, img_height = img.size
                        print(f"[PPTX] Размер: {img_width}x{img_height}")
                        
                        # Создаем слайд
                        end_slide = prs.slides.add_slide(prs.slide_layouts[6])
                        
                        # Конвертируем пиксели в дюймы (предполагаем 96 DPI)
                        img_width_inches = img_width / 96.0
                        img_height_inches = img_height / 96.0
                        
                        # Проверяем, помещается ли изображение на слайд
                        slide_width_inches = prs.slide_width / Inches(1)
                        slide_height_inches = prs.slide_height / Inches(1)
                        
                        # Если изображение больше слайда, масштабируем с сохранением пропорций
                        if img_width_inches > slide_width_inches or img_height_inches > slide_height_inches:
                            scale = min(slide_width_inches / img_width_inches, slide_height_inches / img_height_inches)
                            img_width_inches *= scale
                            img_height_inches *= scale
                            print(f"[PPTX] Изображение масштабировано до: {img_width_inches:.2f}x{img_height_inches:.2f} дюймов")
                        
                        # Центрируем изображение на слайде
                        left = (prs.slide_width - Inches(img_width_inches)) / 2
                        top = (prs.slide_height - Inches(img_height_inches)) / 2
                        
                        end_slide.shapes.add_picture(end_path, left, top, 
                                                    width=Inches(img_width_inches), 
                                                    height=Inches(img_height_inches))
                        print(f"[PPTX] Финальный слайд {end_name} добавлен успешно")
                        end_added = True
                        break
                    except Exception as e:
                        print(f"Ошибка добавления финального слайда {end_path}: {e}")
                        import traceback
                        traceback.print_exc()
            if end_added:
                break
        
        if not end_added:
            print(f"[PPTX] Финальный слайд {end_name} не найден")
    
    # Сохраняем презентацию
    output_path = f"outputs/{project_name.replace(' ', '_')}.pptx"
    os.makedirs('outputs', exist_ok=True)
    prs.save(output_path)
    
    return output_path
